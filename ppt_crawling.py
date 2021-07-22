from pptx import Presentation
import os
import re
import load_phrase as lp
import util

# 훈련용 데이터 저장
x_train = []
y_train = []

# PPT 디렉토리 불러오기
root_dir = os.listdir('CCNA/ppts')

# 단어구 로드
load_phrase = lp.LoadPhrase()
phrases = load_phrase.get_phrases()


for i in range(len(root_dir)):
    prs = Presentation('CCNA/ppts/' + f'{root_dir[i]}')
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.text = util.clear_word(paragraph.text)
                    slide_text.append(paragraph.text)
            elif shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        cell.text = util.clear_word(cell.text)
                        slide_text.append(cell.text)

        slide_text = [text for text in slide_text if len(text) > 1]
        slide_str = ''
        for text in slide_text:
            slide_str += text

        for phrase in phrases:
            slide_str = slide_str.replace(phrase, re.sub(' +', '-', phrase))

        x_train.append(slide_str)
        y_train.append(root_dir[i].split('.')[0])

util.array_to_csv(x_train, 'x_train.csv')
util.array_to_csv(y_train, 'y_train.csv')

util.array_to_csv(set(y_train), 'category.csv')


